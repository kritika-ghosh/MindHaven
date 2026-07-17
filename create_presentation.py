import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

os.makedirs('artifacts', exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
C_BG = RGBColor(11, 12, 17)       # #0b0c11 (Dark background)
C_CARD = RGBColor(19, 20, 31)     # #13141f (Card background)
C_WHITE = RGBColor(255, 255, 255) # Text Primary
C_GRAY = RGBColor(181, 183, 205)  # #b5b7cd (Text Secondary)
C_PURPLE = RGBColor(132, 85, 239) # #8455ef
C_CYAN = RGBColor(0, 242, 254)    # #00f2fe
C_PINK = RGBColor(255, 0, 127)    # #ff007f
C_EMERALD = RGBColor(16, 185, 129)# #10b981
C_RED = RGBColor(239, 68, 68)     # #ef4444

def add_slide_with_background(prs, title_text, slide_num_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Paint background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    
    # Add top banner line (sleek cyan accent)
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(0.12)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = C_CYAN
    banner.line.fill.background()
    
    if title_text:
        # Title box
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
    if slide_num_str:
        # Slide number
        numBox = slide.shapes.add_textbox(Inches(11.8), Inches(0.4), Inches(1.0), Inches(0.8))
        tf = numBox.text_frame
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = slide_num_str
        p.font.name = 'Share Tech Mono'
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_CYAN
        
    return slide

def add_bullet_points(slide, left, top, width, height, points, size_pt=15.5, space_pt=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    for i, pt in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pt
        p.font.name = 'Calibri'
        p.font.size = Pt(size_pt)
        p.font.color.rgb = C_GRAY
        p.space_after = Pt(space_pt)
        
        if "[KEY]" in pt:
            p.text = pt.replace("[KEY]", "")
            p.font.bold = True
            p.font.color.rgb = C_WHITE
    return txBox

# ==========================================
# SLIDE 1: TITLE SLIDE
# ==========================================
slide1 = add_slide_with_background(prs, None, None)

# Title content
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0)
tf.margin_right = Inches(0)
tf.margin_top = Inches(0)
tf.margin_bottom = Inches(0)

# Track badge
p_badge = tf.paragraphs[0]
p_badge.text = "BUILD FOR GOOD HACKATHON  |  SWASTHYA TRACK"
p_badge.font.name = 'Segoe UI'
p_badge.font.size = Pt(11)
p_badge.font.bold = True
p_badge.font.color.rgb = C_CYAN
p_badge.space_after = Pt(16)

# Main Title
p_title = tf.add_paragraph()
p_title.text = "ML-Based Student Burnout\nDetection System"
p_title.font.name = 'Segoe UI'
p_title.font.size = Pt(38)
p_title.font.bold = True
p_title.font.color.rgb = C_WHITE
p_title.space_after = Pt(10)

# Subtitle
p_sub = tf.add_paragraph()
p_sub.text = "A Multimodal Approach Combining Psychometric, Visual & Acoustic Signals"
p_sub.font.name = 'Segoe UI'
p_sub.font.size = Pt(15)
p_sub.font.italic = True
p_sub.font.color.rgb = C_PURPLE
p_sub.space_after = Pt(32)

# Divider line
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(7.0), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = C_PURPLE
line.line.fill.background()

# Author Info
author_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(7.5), Inches(2.2))
tf_auth = author_box.text_frame
tf_auth.word_wrap = True
tf_auth.margin_left = Inches(0)
tf_auth.margin_right = Inches(0)
tf_auth.margin_top = Inches(0)
tf_auth.margin_bottom = Inches(0)

p_authors = tf_auth.paragraphs[0]
p_authors.text = "Presenter: Kritika Ghosh"
p_authors.font.name = 'Calibri'
p_authors.font.size = Pt(15)
p_authors.font.bold = True
p_authors.font.color.rgb = C_WHITE
p_authors.space_after = Pt(8)

p_school = tf_auth.add_paragraph()
p_school.text = "School of Computing Science Engineering & AI | VIT Bhopal University\nSeptember 2025 Submission"
p_school.font.name = 'Calibri'
p_school.font.size = Pt(11)
p_school.font.color.rgb = C_GRAY

# Right side decorative orb
orb_outer = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(1.8), Inches(3.8), Inches(3.8))
orb_outer.fill.solid()
orb_outer.fill.fore_color.rgb = RGBColor(15, 23, 42)
orb_outer.line.color.rgb = C_PURPLE
orb_outer.line.width = Pt(2)

orb_inner = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(2.1), Inches(3.2), Inches(3.2))
orb_inner.fill.solid()
orb_inner.fill.fore_color.rgb = RGBColor(22, 28, 45)
orb_inner.line.color.rgb = C_CYAN
orb_inner.line.width = Pt(1.5)

orb_text = slide1.shapes.add_textbox(Inches(9.2), Inches(3.1), Inches(3.0), Inches(1.2))
tf_orb = orb_text.text_frame
tf_orb.word_wrap = True
tf_orb.margin_left = Inches(0)
tf_orb.margin_right = Inches(0)
tf_orb.margin_top = Inches(0)
tf_orb.margin_bottom = Inches(0)

p_orb = tf_orb.paragraphs[0]
p_orb.alignment = PP_ALIGN.CENTER
p_orb.text = "SWASTHYA\nTRACK"
p_orb.font.name = 'Segoe UI'
p_orb.font.size = Pt(20)
p_orb.font.bold = True
p_orb.font.color.rgb = C_WHITE

p_orb_sub = tf_orb.add_paragraph()
p_orb_sub.alignment = PP_ALIGN.CENTER
p_orb_sub.text = "Mental Health for Youth"
p_orb_sub.font.name = 'Calibri'
p_orb_sub.font.size = Pt(11)
p_orb_sub.font.color.rgb = C_GRAY


# ==========================================
# SLIDE 2: INTRODUCTION
# ==========================================
slide2 = add_slide_with_background(prs, "Introduction — Youth Mental Health Challenge", "01")

points2 = [
    "[KEY]What is Mental Burnout?",
    "A state of chronic academic stress resulting from prolonged study pressures, leading to:",
    "  • Emotional exhaustion & severe psychological fatigue",
    "  • Marked drop in academic productivity & self-efficacy",
    "  • Depersonalisation, detachment, and social withdrawal",
    "",
    "[KEY]Why Focus on College Students?",
    "Students face unique, cumulative pressures:",
    "  • Academic workload coupled with hostel adjustments",
    "  • Burnout goes unrecognized until reaching a clinical crisis stage",
    "  • Severe lack of accessible, private, non-clinical wellness check-ins"
]
# No negative space because this fills the slide
add_bullet_points(slide2, Inches(0.6), Inches(1.4), Inches(6.2), Inches(5.3), points2, size_pt=14, space_pt=6)

img_path_2 = 'artifacts/diagnostic_terminal_quiz.png'
if os.path.exists(img_path_2):
    slide2.shapes.add_picture(img_path_2, Inches(7.1), Inches(1.6), width=Inches(5.6), height=Inches(4.2))


# ==========================================
# SLIDE 3: LITERATURE & RESEARCH GAP
# ==========================================
slide3 = add_slide_with_background(prs, "Literature Review & Research Gap", "02")

# Comparison Cards are expanded slightly in height to cover negative space
lim_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4))
lim_card.fill.solid()
lim_card.fill.fore_color.rgb = C_CARD
lim_card.line.color.rgb = C_RED
lim_card.line.width = Pt(1.5)

tf_lim = lim_card.text_frame
tf_lim.word_wrap = True
tf_lim.margin_left = Inches(0.3)
tf_lim.margin_right = Inches(0.3)
tf_lim.margin_top = Inches(0.3)
tf_lim.margin_bottom = Inches(0.3)

p = tf_lim.paragraphs[0]
p.text = "❌ Existing Approaches — Limitations"
p.font.name = 'Segoe UI'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_RED
p.space_after = Pt(18)

lims = [
    "• Relies solely on subjective self-report questionnaires (MBI, etc.) which are highly prone to individual response bias.",
    "• Single-modality designs (only surveys OR only static video) fail to capture real-time physiological and behavioral fluctuations.",
    "• Lack of student-specific behavioral inputs (like voice tone, facial tension shifts during academic tasks).",
    "• High accessibility barriers — typical systems require lab equipment, expensive physical sensors, or clinical visits."
]
for pt in lims:
    p = tf_lim.add_paragraph()
    p.text = pt
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = C_GRAY
    p.space_after = Pt(12)

our_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.4))
our_card.fill.solid()
our_card.fill.fore_color.rgb = C_CARD
our_card.line.color.rgb = C_EMERALD
our_card.line.width = Pt(1.5)

tf_our = our_card.text_frame
tf_our.word_wrap = True
tf_our.margin_left = Inches(0.3)
tf_our.margin_right = Inches(0.3)
tf_our.margin_top = Inches(0.3)
tf_our.margin_bottom = Inches(0.3)

p = tf_our.paragraphs[0]
p.text = "✅ Our Approach — Fills the Gap"
p.font.name = 'Segoe UI'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_EMERALD
p.space_after = Pt(18)

ours = [
    "• Fuses three separate streams: Psychometric Survey + Facial Biometrics (EAR/MAR) + Voice VADER Sentiment.",
    "• Custom-collected dataset of 328 student participants containing validated chronic burnout labels.",
    "• Non-intrusive edge tracking: live webcam scans (Mediapipe) and audio recording directly inside the browser.",
    "• Stepped-care output: yields an objective 0–100 score with 5 severity tiers and a seeded CBT Chatbot controller."
]
for pt in ours:
    p = tf_our.add_paragraph()
    p.text = pt
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = C_GRAY
    p.space_after = Pt(12)


# ==========================================
# SLIDE 4: PROPOSED METHODOLOGY
# ==========================================
slide4 = add_slide_with_background(prs, "Proposed Methodology & Data Flow", "03")

points4 = [
    "[KEY]1. Multimodal Data Collection",
    "  • Questionnaire: 5-question Likert scale.",
    "  • Visual: 20-sec webcam recording (Mediapipe landmarks).",
    "  • Acoustic: 10-sec voice sentiment (microphone).",
    "",
    "[KEY]2. Preprocessing & Fusion Pipeline",
    "  • Consolidated into a unified 18-feature vector.",
    "  • StandardScaler normalization applied.",
    "  • Stratified train/test split (80% / 20%).",
    "",
    "[KEY]3. ML Classification & Scoring",
    "  • CatBoost Regressor maps features to the continuous target.",
    "  • Output: Objective 0–100 Burnout Score.",
    "  • Mapped to 5 severity levels for personalized coping plans."
]
# Spread out vertically to reduce negative space
add_bullet_points(slide4, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.3), points4, size_pt=15.5, space_pt=10)

img_path_4 = 'artifacts/system_architecture.png'
if os.path.exists(img_path_4):
    # Centered vertically on the right
    slide4.shapes.add_picture(img_path_4, Inches(6.9), Inches(2.0), width=Inches(5.8), height=Inches(3.4))


# ==========================================
# SLIDE 5: EDGE FEATURE EXTRACTION
# ==========================================
slide5 = add_slide_with_background(prs, "Edge Feature Extraction (No-Hardware Biometrics)", "04")

points5 = [
    "[KEY]Webcam Facial Tracking (Mediapipe & OpenCV)",
    "  • Eye Aspect Ratio (EAR) captures blink frequencies, flutter, and drowsiness anomalies.",
    "  • Mouth Aspect Ratio (MAR) measures jaw tension, stress yawning, and structural tightness.",
    "  • Real-time facial expressions categorized into positive, neutral, and negative state percentages.",
    "",
    "[KEY]Microphone Voice Analysis (Speech API & VADER)",
    "  • Speech-to-Text transcripter decodes student speech parameters on edge.",
    "  • VADER Lexicon engine computes compound sentiment scores (valency range: -1.0 to +1.0) of spoken statements."
]
# Spaced appropriately above the bottom screenshots
add_bullet_points(slide5, Inches(0.6), Inches(1.3), Inches(12.1), Inches(2.3), points5, size_pt=13, space_pt=3)

img_face = 'artifacts/diagnostic_terminal_face.png'
img_voice = 'artifacts/diagnostic_terminal_voice.png'
if os.path.exists(img_face) and os.path.exists(img_voice):
    slide5.shapes.add_picture(img_face, Inches(0.8), Inches(4.1), width=Inches(5.4), height=Inches(3.0))
    slide5.shapes.add_picture(img_voice, Inches(7.0), Inches(4.1), width=Inches(5.4), height=Inches(3.0))


# ==========================================
# SLIDE 6: DATASET DESCRIPTION
# ==========================================
slide6 = add_slide_with_background(prs, "Dataset Description & Modalities", "05")

# Stats Cards (Fixed fourth card to match website's 95.76% R2 instead of old 78% accuracy)
stats = [
    ("328", "Total Participants", 0.6, C_CYAN),
    ("3", "Data Modalities", 2.2, C_PURPLE),
    ("5", "Burnout Classes", 3.8, C_PINK),
    ("95.76%", "Model R² Score", 5.4, C_EMERALD)
]
for val, title, left, col in stats:
    card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.4), Inches(1.4), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD
    card.line.color.rgb = col
    card.line.width = Pt(1.5)
    
    tf_stat = card.text_frame
    tf_stat.word_wrap = True
    tf_stat.margin_left = Inches(0)
    tf_stat.margin_right = Inches(0)
    tf_stat.margin_top = Inches(0.1)
    tf_stat.margin_bottom = Inches(0)
    
    p = tf_stat.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = val
    p.font.name = 'Segoe UI'
    # Adjusted font size slightly for 95.76% string fit
    p.font.size = Pt(22) if "%" in val and len(val)>4 else Pt(26)
    p.font.bold = True
    p.font.color.rgb = col
    
    p2 = tf_stat.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = title
    p2.font.name = 'Calibri'
    p2.font.size = Pt(10)
    p2.font.bold = True
    p2.font.color.rgb = C_GRAY

points6 = [
    "[KEY]Dataset Features Summary",
    "  • Survey (Q1–Q5): Likert scale metrics capturing exhaustion, detachment, avoidance, academic pride, and outcome recognition.",
    "  • Survey (Q6 / Label): Single-item validated chronic burnout scale mapping to 5 classes (0.0 to 4.0).",
    "  • Facial Features: Eye fatigue (Avg EAR), jaw tension (Avg MAR), and expression distribution (Positive, Neutral, Negative %).",
    "  • Voice Features: Spoken valence sentiment scores via VADER lexicon-based sentiment analysis."
]
add_bullet_points(slide6, Inches(0.6), Inches(2.9), Inches(6.2), Inches(4.0), points6, size_pt=13.5, space_pt=8)

# Table on the right
rows, cols = 7, 3
left_t, top_t, width_t, height_t = Inches(7.1), Inches(1.4), Inches(5.6), Inches(5.2)
table_shape = slide6.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
table = table_shape.table

table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(1.4)
table.columns[2].width = Inches(3.2)

headers = ["Modality", "Feature", "Description"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PURPLE
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Segoe UI'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

data_rows = [
    ["Survey", "Q1–Q5", "Likert scale (0–4) on burnout symptomatology"],
    ["Survey", "Q6 (Label)", "Validated single-item chronic burnout index"],
    ["Facial", "Avg EAR", "Eye Aspect Ratio (fatigue/blink rate)"],
    ["Facial", "Avg MAR", "Mouth Aspect Ratio (tension/jaw tightness)"],
    ["Facial", "Emotion %", "Positive / Neutral / Negative face percentage"],
    ["Voice", "VADER", "Compound spoken valence sentiment scores"]
]

for row_idx, row_data in enumerate(data_rows):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_CARD
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = C_GRAY


# ==========================================
# SLIDE 7: MODEL EVALUATION & RESULTS
# ==========================================
slide7 = add_slide_with_background(prs, "Experimental Results: Model Comparison", "06")

points7 = [
    "[KEY]Regression Model Benchmarks (5-Fold CV)",
    "We compared multiple regression architectures to predict the continuous burnout index on our student dataset:",
    "",
    "  • [KEY]Support Vector Regressor (SVR)",
    "    Mean R² Score: [KEY]92.95%  |  Mean RMSE: [KEY]0.2227  |  Mean MAE: [KEY]0.1719",
    "",
    "  • [KEY]Random Forest Regressor",
    "    Mean R² Score: [KEY]95.32%  |  Mean RMSE: [KEY]0.1805  |  Mean MAE: [KEY]0.1444",
    "",
    "  • [KEY]CatBoost Regressor (Ours — Final Model)",
    "    Mean R² Score: [KEY]95.76%  |  Mean RMSE: [KEY]0.1727  |  Mean MAE: [KEY]0.1378",
    "    Consistently achieves the highest continuous prediction performance with standard error below 0.18."
]
# Spread out bullet lines to fill slide height
add_bullet_points(slide7, Inches(0.6), Inches(1.4), Inches(6.4), Inches(5.3), points7, size_pt=14, space_pt=10)

img_path_7 = 'artifacts/cross_validation_card.png'
if os.path.exists(img_path_7):
    # Positioned nicely in the middle right
    slide7.shapes.add_picture(img_path_7, Inches(7.2), Inches(2.4), width=Inches(5.4), height=Inches(2.7))


# ==========================================
# SLIDE 8: FEATURE IMPORTANCE
# ==========================================
slide8 = add_slide_with_background(prs, "Explainable AI: Feature Importance Analysis", "07")

points8 = [
    "[KEY]Predictive Weights Validation (Regression)",
    "CatBoost feature weighting lists the contribution of inputs to the continuous burnout score target:",
    "",
    "  • [KEY]Consolidated Survey (Survey_Sum): Ranks #1 (weight: 77.70%).",
    "    Directly represents the baseline questionnaire severity score.",
    "",
    "  • [KEY]Biometric & Spoken Feature Refinement (22.30% cumulative):",
    "    Fine-tunes the baseline score to account for objective biological fatigue:",
    "      - Progress Pride (Q4_inv): [KEY]3.85% (Ranks #2)",
    "      - Meaningful Outcomes (Q5_inv): [KEY]3.41% (Ranks #3)",
    "      - Social Avoidance (Q3): [KEY]2.22% (Ranks #4)",
    "      - Vocal Negative Tone (Sentiment_Neg): [KEY]1.95% (Ranks #6)",
    "      - Mouth Aspect Ratio (Avg_MAR): [KEY]1.29% (Ranks #7)"
]
# Spaced to fill height
add_bullet_points(slide8, Inches(0.6), Inches(1.4), Inches(6.4), Inches(5.3), points8, size_pt=13, space_pt=8)

img_path_8 = 'artifacts/feature_importance.png'
if os.path.exists(img_path_8):
    # Scaled and centered to fill vertical negative space on the right
    slide8.shapes.add_picture(img_path_8, Inches(7.2), Inches(1.6), width=Inches(5.4), height=Inches(4.2))


# ==========================================
# SLIDE 9: NEW CHATBOT SLIDE (NEW!)
# ==========================================
slide9 = add_slide_with_background(prs, "MindHaven-CBT: The Empathetic AI Companion", "08")

points9 = [
    "[KEY]Fine-Tuned Cognitive Therapy Core",
    "The client portal integrates a context-aware chatbot specially trained to serve student mental wellness needs offline:",
    "",
    "  • [KEY]Transfer Learning & SFT:",
    "    Qwen base LLM fine-tuned offline on real-world therapist responses (CounselChat + Mental Health FAQ) using QLoRA adapters.",
    "",
    "  • [KEY]CBT Empathetic Guidelines:",
    "    Identifies cognitive distortions (catastrophizing, academic imposter syndrome) to guide the student toward structured reflection.",
    "",
    "  • [KEY]Safety Guardrail Layer:",
    "    Includes automatic distress detection to flag severe text sequences and instantly route users to university wellness contacts.",
    "",
    "  • [KEY]Stepped-Care Micro-Plans:",
    "    Generates 7-day personalized wellness tasks like breathing scripts and workspace boundary boundaries."
]
add_bullet_points(slide9, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.3), points9, size_pt=13.5, space_pt=8)

img_path_9 = 'artifacts/telemetry_insights.png'
if os.path.exists(img_path_9):
    # Centered and scaled to fill right side
    slide9.shapes.add_picture(img_path_9, Inches(6.9), Inches(1.6), width=Inches(5.8), height=Inches(4.35))


# ==========================================
# SLIDE 10: SWASTHYA APPLICATION (INSIGHTS)
# ==========================================
slide10 = add_slide_with_background(prs, "Swasthya Integration: Biometric Context Seeding", "09")

points10 = [
    "[KEY]Empathetic Context Injection",
    "Traditional wellness apps treat chatbots as disconnected conversational loops. MindHaven seeds context before text begins:",
    "",
    "  • [KEY]Automated Prompt Seeding:",
    "    Prior to launching, the chatbot's system instructions are dynamically seeded with the latest assessment parameters (burnout score, EAR/MAR fatigue metrics, expression distribution, and voice sentiment).",
    "",
    "  • [KEY]Biometric Self-Awareness:",
    "    The AI companion knows exactly why the user scored a 3.2/4.0 (e.g. noticing eye-fatigue and negative tone) and references this to deliver structured Cognitive Behavioral Therapy suggestions.",
    "",
    "  • [KEY]Stepped-Care Exercises:",
    "    Integrates directly with the Pranayama Studio breathing orb to guide students through deep stress recovery cycles."
]
# Spread points to reduce bottom space
add_bullet_points(slide10, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.3), points10, size_pt=14, space_pt=10)

img_path_pran = 'artifacts/pranayama_studio.png'
if os.path.exists(img_path_pran):
    # Centered and scaled on the right
    slide10.shapes.add_picture(img_path_pran, Inches(6.9), Inches(1.6), width=Inches(5.8), height=Inches(4.35))


# ==========================================
# SLIDE 11: CONCLUSION & ROADMAP
# ==========================================
slide11 = add_slide_with_background(prs, "Conclusion & Swasthya Future Roadmap", "10")

points11 = [
    "[KEY]Presentation Key Takeaways",
    "  • [KEY]Multimodal Fusion: First system to combine psychometric, visual (EAR/MAR), and voice (VADER) in one non-intrusive student wellness tool.",
    "  • [KEY]Proven Results: CatBoost model yields 78% classification accuracy, and continuous regression model achieves 95.76% R² variance mapping.",
    "  • [KEY]Actionable Care: Incorporates glassmorphic CBT companion and expanding breathing orb for immediate stress reset.",
    "",
    "[KEY]Deployment Roadmap",
    "  • [KEY]Phase 1: Retraining: Expand the 328-participant dataset across demographics.",
    "  • [KEY]Phase 2: Physiology: Integrate heart rate variance (HRV via photoplethysmography).",
    "  • [KEY]Phase 3: Mobile Apps: Launch real-time edge wellness log applications.",
    "  • [KEY]Phase 4: University Pilot: Implement automated student support trials."
]
add_bullet_points(slide11, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.3), points11, size_pt=14, space_pt=10)

img_path_11 = 'artifacts/telemetry_dashboard.png'
if os.path.exists(img_path_11):
    slide11.shapes.add_picture(img_path_11, Inches(6.9), Inches(1.6), width=Inches(5.8), height=Inches(4.35))


# ==========================================
# SLIDE 12: THANK YOU
# ==========================================
slide12 = add_slide_with_background(prs, "Thank You", "11")

center_card = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.8))
center_card.fill.solid()
center_card.fill.fore_color.rgb = C_CARD
center_card.line.color.rgb = C_CYAN
center_card.line.width = Pt(2)

tf_thank = center_card.text_frame
tf_thank.word_wrap = True
tf_thank.margin_left = Inches(0.5)
tf_thank.margin_right = Inches(0.5)
tf_thank.margin_top = Inches(0.5)
tf_thank.margin_bottom = Inches(0.5)

p_th = tf_thank.paragraphs[0]
p_th.alignment = PP_ALIGN.CENTER
p_th.text = "Questions & Discussion Welcome"
p_th.font.name = 'Segoe UI'
p_th.font.size = Pt(28)
p_th.font.bold = True
p_th.font.color.rgb = C_WHITE
p_th.space_after = Pt(24)

p_presents = tf_thank.add_paragraph()
p_presents.alignment = PP_ALIGN.CENTER
p_presents.text = "Presenter: Kritika Ghosh"
p_presents.font.name = 'Segoe UI'
p_presents.font.size = Pt(18)
p_presents.font.bold = True
p_presents.font.color.rgb = C_PURPLE
p_presents.space_after = Pt(18)

p_github = tf_thank.add_paragraph()
p_github.alignment = PP_ALIGN.CENTER
p_github.text = "GitHub Repository: github.com/kritika-ghosh/Mental-Burnout-Detection-System"
p_github.font.name = 'Share Tech Mono'
p_github.font.size = Pt(14)
p_github.font.color.rgb = C_CYAN
p_github.space_after = Pt(8)

p_demo = tf_thank.add_paragraph()
p_demo.alignment = PP_ALIGN.CENTER
p_demo.text = "Live Demo Space: huggingface.co/spaces/project-exhibition/Burnout-detection"
p_demo.font.name = 'Share Tech Mono'
p_demo.font.size = Pt(14)
p_demo.font.color.rgb = C_CYAN
p_demo.space_after = Pt(28)

p_footer = tf_thank.add_paragraph()
p_footer.alignment = PP_ALIGN.CENTER
p_footer.text = "SWASTHYA Track Submission  |  Build for Good Hackathon  |  School of CSE & AI, VIT Bhopal"
p_footer.font.name = 'Calibri'
p_footer.font.size = Pt(11)
p_footer.font.color.rgb = C_GRAY

saved = False
for idx in range(1, 20):
    suffix = "" if idx == 1 else f"_v{idx}"
    target_path = f"artifacts/Swasthya_MindHaven_Submission{suffix}.pptx"
    try:
        prs.save(target_path)
        print(f"SWASHTYA PowerPoint submission slide deck successfully updated and saved to {target_path}!")
        saved = True
        # Keep track of the actual saved path to print it out
        actual_saved_path = target_path
        break
    except PermissionError:
        continue
if not saved:
    print("Error: Could not save presentation. All filename increments (v1-v19) were locked.")
